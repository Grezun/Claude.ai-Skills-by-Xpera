#!/usr/bin/env python3
"""Build a branded .pptx deck from corporate markdown + a brand template.

Pipeline:
  1. Parse markdown → structured outline (md_to_outline.py)
  2. Open template → inherit master, theme, and color palette
  3. Clear template placeholder slides
  4. Add slides in outline order, each linked to its layout
  5. Fill title and body placeholders from outline content
  6. Save bundle

Usage:
  python scripts/build_deck.py input.md template.pptx output.pptx
  python scripts/build_deck.py input.md template.pptx output.pptx --dry-run
"""
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.oxml.ns import qn
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# Add scripts/ to path so md_to_outline is importable
sys.path.insert(0, str(Path(__file__).parent))
from md_to_outline import parse_markdown


# Fragments used to identify layouts by name (case-insensitive substring match)
# Keys map to slide types emitted by md_to_outline.py
LAYOUT_HINTS = {
    "cover":    ["title slide", "cover", "title only"],
    "section":  ["section header", "section", "divider"],
    "content":  ["title and content", "content", "body"],
    "callout":  ["quote", "callout", "blank", "statement"],
    "data":     ["two content", "comparison", "two column", "data"],
    "image":    ["picture with caption", "picture", "image", "blank"],
}

# Fallback priority when no hint matches
FALLBACK_ORDER = ["content", "cover", "section"]


def find_layout(master, slide_type: str):
    """Return the best-matching slide layout object for a given slide type."""
    hints = LAYOUT_HINTS.get(slide_type, LAYOUT_HINTS["content"])
    layouts = master.slide_layouts

    for hint in hints:
        for layout in layouts:
            if hint in layout.name.lower():
                return layout

    # Generic fallback — use second layout (index 1) or first
    return layouts[min(1, len(layouts) - 1)]


def clear_template_slides(prs: Presentation) -> None:
    """Remove all placeholder slides added by the template file itself."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)


def fill_placeholder(slide, idx: int, text: str) -> bool:
    """Write text into a placeholder by index (0=title, 1=body). Returns True if found."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return True
    return False


def fill_body_bullets(slide, bullets: list[str]) -> None:
    """Write bullet list into body placeholder (idx=1), one paragraph per bullet."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
            return


def add_slide(prs: Presentation, layout, title: str, bullets: list[str], dry_run: bool) -> None:
    if dry_run:
        bullet_preview = f" ({len(bullets)} bullets)" if bullets else ""
        print(f"    + [{layout.name}] {title!r}{bullet_preview}")
        return

    slide = prs.slides.add_slide(layout)
    fill_placeholder(slide, 0, title)
    if bullets:
        fill_body_bullets(slide, bullets)


def build_deck(md_path: str, template_path: str, output_path: str, dry_run: bool = False) -> int:
    outline = parse_markdown(md_path)
    prs = Presentation(template_path)
    master = prs.slide_master

    if not dry_run:
        clear_template_slides(prs)

    slide_count = 0

    # Cover slide
    cover_layout = find_layout(master, "cover")
    subtitle = outline.get("subtitle", "")
    if dry_run:
        print(f"DECK: {outline['title']!r}")
        print(f"  + [{cover_layout.name}] {outline['title']!r} (cover)")
    else:
        slide = prs.slides.add_slide(cover_layout)
        fill_placeholder(slide, 0, outline["title"])
        if subtitle:
            fill_placeholder(slide, 1, subtitle)
    slide_count += 1

    for section in outline["sections"]:
        # Section divider
        if section["title"]:
            sec_layout = find_layout(master, "section")
            if dry_run:
                print(f"\n  SECTION: {section['title']!r}")
                print(f"    + [{sec_layout.name}] {section['title']!r} (divider)")
            else:
                slide = prs.slides.add_slide(sec_layout)
                fill_placeholder(slide, 0, section["title"])
            slide_count += 1
        elif dry_run:
            print(f"\n  SECTION: (untitled)")

        # Content slides
        for s in section["slides"]:
            slide_type = s.get("type", "content")
            layout = find_layout(master, slide_type)
            bullets = s.get("bullets", [])
            add_slide(prs, layout, s["title"], bullets, dry_run)
            slide_count += 1

    if not dry_run:
        prs.save(output_path)
        print(f"Saved {output_path} ({slide_count} slides)")
    else:
        print(f"\nDry run complete — would produce {slide_count} slides")

    return slide_count


def main():
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    dry_run = "--dry-run" in sys.argv

    if len(args) < 3:
        print("Usage: build_deck.py <input.md> <template.pptx> <output.pptx> [--dry-run]", file=sys.stderr)
        sys.exit(1)

    build_deck(args[0], args[1], args[2], dry_run=dry_run)


if __name__ == "__main__":
    main()
